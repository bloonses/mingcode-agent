"""Office 文档读写工具集。

包含 4 个工具：
- WordTool: 读写 .docx（依赖 python-docx）
- PdfTool: 读 .pdf（pypdf）+ 写 .pdf（reportlab）
- ExcelTool: 读写 .xlsx（openpyxl）
- PptTool: 读写 .pptx（python-pptx）

所有 Office 库均为可选导入，缺失时返回友好错误字符串而非抛 ImportError。
"""
import json as _json
import os
from typing import Optional, List, Any

from .base import BaseTool


# ============== 可选依赖导入辅助函数 ==============

def _import_docx():
    try:
        from docx import Document
        return Document
    except ImportError:
        return None


def _import_pypdf():
    try:
        from pypdf import PdfReader
        return PdfReader
    except ImportError:
        return None


def _import_reportlab_canvas():
    try:
        from reportlab.pdfgen import canvas
        return canvas
    except ImportError:
        return None


def _import_openpyxl():
    try:
        import openpyxl
        return openpyxl
    except ImportError:
        return None


def _import_pptx():
    try:
        from pptx import Presentation
        return Presentation
    except ImportError:
        return None


# ============== WordTool ==============

class WordTool(BaseTool):
    name = "word"
    description = (
        "读写 Microsoft Word .docx 文档。支持 action："
        "read（读取所有段落文本，可选 max_chars 限制返回字符数，默认 8000）/ "
        "write（按 \\n 分段写入新文档，需提供 content）。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "write"],
                "description": "操作类型"
            },
            "path": {
                "type": "string",
                "description": "文件路径"
            },
            "content": {
                "type": "string",
                "description": "action=write 时必填，要写入的文本（段落用 \\n 分隔）"
            },
            "max_chars": {
                "type": "integer",
                "description": "action=read 时可选，最大返回字符数，默认 8000"
            }
        },
        "required": ["action", "path"]
    }

    def execute(self, **kwargs) -> str:
        action = (kwargs.get("action") or "").strip().lower()
        path = kwargs.get("path")

        if not path:
            return "Error: path is required"
        if action not in ("read", "write"):
            return f"Error: unknown action '{action}'. Supported: read/write"

        Document = _import_docx()
        if Document is None:
            return "Error: python-docx 未安装。请运行: pip install python-docx"

        if action == "read":
            return self._read(Document, path, kwargs.get("max_chars"))
        return self._write(Document, path, kwargs.get("content"))

    def _read(self, Document, path: str, max_chars) -> str:
        if not os.path.exists(path):
            return f"Error: file not found: {path}"
        try:
            doc = Document(path)
        except Exception as e:
            return f"Error: failed to open docx: {e}"

        limit = max_chars if isinstance(max_chars, int) and max_chars > 0 else 8000
        parts: List[str] = []
        total = 0
        for para in doc.paragraphs:
            text = para.text
            if total + len(text) > limit:
                remain = limit - total
                if remain > 0:
                    parts.append(text[:remain])
                parts.append("\n... [已截断，超过最大字符数限制]")
                break
            parts.append(text)
            total += len(text)
        return "\n".join(parts)

    def _write(self, Document, path: str, content) -> str:
        if content is None:
            return "Error: content is required for write action"
        content = str(content)

        parent_dir = os.path.dirname(path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)

        try:
            doc = Document()
            for line in content.split("\n"):
                doc.add_paragraph(line)
            doc.save(path)
        except Exception as e:
            return f"Error: failed to write docx: {e}"
        return f"Word 文档已成功写入：{path}"


# ============== PdfTool ==============

class PdfTool(BaseTool):
    name = "pdf"
    description = (
        "读写 PDF 文档。支持 action："
        "read（提取每页文本，格式为 [Page N]\\n<text>，可选 max_pages 限制，默认 50）/ "
        "write（按 \\n 分行写入纯文本 PDF，需提供 content）。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "write"],
                "description": "操作类型"
            },
            "path": {
                "type": "string",
                "description": "文件路径"
            },
            "content": {
                "type": "string",
                "description": "action=write 时必填，要写入的纯文本内容"
            },
            "max_pages": {
                "type": "integer",
                "description": "action=read 时可选，最大读取页数，默认 50"
            }
        },
        "required": ["action", "path"]
    }

    def execute(self, **kwargs) -> str:
        action = (kwargs.get("action") or "").strip().lower()
        path = kwargs.get("path")

        if not path:
            return "Error: path is required"
        if action not in ("read", "write"):
            return f"Error: unknown action '{action}'. Supported: read/write"

        if action == "read":
            PdfReader = _import_pypdf()
            if PdfReader is None:
                return "Error: pypdf 未安装。请运行: pip install pypdf"
            return self._read(PdfReader, path, kwargs.get("max_pages"))

        canvas = _import_reportlab_canvas()
        if canvas is None:
            return "Error: reportlab 未安装。请运行: pip install reportlab"
        return self._write(canvas, path, kwargs.get("content"))

    def _read(self, PdfReader, path: str, max_pages) -> str:
        if not os.path.exists(path):
            return f"Error: file not found: {path}"
        try:
            reader = PdfReader(path)
        except Exception as e:
            return f"Error: failed to open pdf: {e}"

        limit = max_pages if isinstance(max_pages, int) and max_pages > 0 else 50
        pages = reader.pages[:limit]
        parts: List[str] = []
        for i, page in enumerate(pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            parts.append(f"[Page {i}]\n{text}")
        if len(reader.pages) > limit:
            parts.append(f"... [已截断，共 {len(reader.pages)} 页，仅显示前 {limit} 页]")
        return "\n\n".join(parts)

    def _write(self, canvas, path: str, content) -> str:
        if content is None:
            return "Error: content is required for write action"
        content = str(content)

        parent_dir = os.path.dirname(path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)

        try:
            c = canvas.Canvas(path)
            # A4 页面尺寸（点）：宽 595.28，高 841.89
            width, height = 595.28, 841.89
            margin = 50
            line_height = 14
            x = margin
            y = height - margin

            lines = content.split("\n")
            for line in lines:
                if y < margin:
                    c.showPage()
                    y = height - margin
                c.drawString(x, y, line)
                y -= line_height
            c.save()
        except Exception as e:
            return f"Error: failed to write pdf: {e}"
        return f"PDF 文档已成功写入：{path}"


# ============== ExcelTool ==============

class ExcelTool(BaseTool):
    name = "excel"
    description = (
        "读写 Microsoft Excel .xlsx 文件。支持 action："
        "read（读取指定 sheet，返回 Markdown 表格，可选 max_rows 限制，默认 100）/ "
        "write（写入二维数组 data，如 [[\"姓名\",\"年龄\"],[\"张三\",20]]）。"
        "可选 sheet 参数指定工作表名称，默认 Sheet1。"
    )
    parameters = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "write"],
                "description": "操作类型"
            },
            "path": {
                "type": "string",
                "description": "文件路径"
            },
            "sheet": {
                "type": "string",
                "description": "工作表名称，默认 Sheet1"
            },
            "data": {
                "type": "string",
                "description": "action=write 时必填，二维数组的 JSON 字符串，如 [[\"姓名\",\"年龄\"],[\"张三\",20]]"
            },
            "max_rows": {
                "type": "integer",
                "description": "action=read 时可选，最大读取行数，默认 100"
            }
        },
        "required": ["action", "path"]
    }

    def execute(self, **kwargs) -> str:
        action = (kwargs.get("action") or "").strip().lower()
        path = kwargs.get("path")

        if not path:
            return "Error: path is required"
        if action not in ("read", "write"):
            return f"Error: unknown action '{action}'. Supported: read/write"

        openpyxl = _import_openpyxl()
        if openpyxl is None:
            return "Error: openpyxl 未安装。请运行: pip install openpyxl"

        sheet = kwargs.get("sheet") or "Sheet1"

        if action == "read":
            return self._read(openpyxl, path, sheet, kwargs.get("max_rows"))
        return self._write(openpyxl, path, sheet, kwargs.get("data"))

    def _read(self, openpyxl, path: str, sheet: str, max_rows) -> str:
        if not os.path.exists(path):
            return f"Error: file not found: {path}"
        try:
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        except Exception as e:
            return f"Error: failed to open xlsx: {e}"

        if sheet not in wb.sheetnames:
            wb.close()
            return f"Error: sheet '{sheet}' not found. Available: {wb.sheetnames}"

        ws = wb[sheet]
        limit = max_rows if isinstance(max_rows, int) and max_rows > 0 else 100

        rows: List[List[str]] = []
        count = 0
        for row in ws.iter_rows(values_only=True):
            if count >= limit:
                break
            cells = ["" if v is None else str(v) for v in row]
            rows.append(cells)
            count += 1
        wb.close()

        if not rows:
            return "(空表)"

        # 转 Markdown 表格
        # 表头
        header = rows[0]
        col_count = len(header)
        md = []
        md.append("| " + " | ".join(header) + " |")
        md.append("| " + " | ".join(["---"] * col_count) + " |")
        for r in rows[1:]:
            # 补齐列数
            padded = r + [""] * (col_count - len(r))
            md.append("| " + " | ".join(padded[:col_count]) + " |")
        if count >= limit:
            md.append(f"\n... [已截断，仅显示前 {limit} 行]")
        return "\n".join(md)

    def _write(self, openpyxl, path: str, sheet: str, data) -> str:
        if data is None:
            return "Error: data is required for write action"
        # data 是 JSON 字符串
        try:
            parsed = _json.loads(data) if isinstance(data, str) else data
        except Exception as e:
            return f"Error: failed to parse data as JSON: {e}"

        if not isinstance(parsed, list) or not all(isinstance(r, list) for r in parsed):
            return "Error: data must be a 2D array (list of lists)"

        parent_dir = os.path.dirname(path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)

        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet
            for row in parsed:
                ws.append(row)
            wb.save(path)
            wb.close()
        except Exception as e:
            return f"Error: failed to write xlsx: {e}"
        return f"Excel 文档已成功写入：{path}"


# ============== PptTool ==============

class PptTool(BaseTool):
    description = (
        "读写 Microsoft PowerPoint .pptx 演示文稿。支持 action："
        "read（遍历每张幻灯片，提取标题和所有文本框，格式为 [Slide N] Title: ...\\n<text>，可选 max_slides 限制，默认 50）/ "
        "write（按 slides 数组创建幻灯片，每项 {\"title\": \"...\", \"content\": \"...\"}，content 多行用 \\n 分隔）。"
    )
    name = "ppt"
    parameters = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["read", "write"],
                "description": "操作类型"
            },
            "path": {
                "type": "string",
                "description": "文件路径"
            },
            "slides": {
                "type": "string",
                "description": "action=write 时必填，幻灯片数组的 JSON 字符串，每项 {\"title\": \"...\", \"content\": \"...\"}"
            },
            "max_slides": {
                "type": "integer",
                "description": "action=read 时可选，最大读取幻灯片数，默认 50"
            }
        },
        "required": ["action", "path"]
    }

    def execute(self, **kwargs) -> str:
        action = (kwargs.get("action") or "").strip().lower()
        path = kwargs.get("path")

        if not path:
            return "Error: path is required"
        if action not in ("read", "write"):
            return f"Error: unknown action '{action}'. Supported: read/write"

        Presentation = _import_pptx()
        if Presentation is None:
            return "Error: python-pptx 未安装。请运行: pip install python-pptx"

        if action == "read":
            return self._read(Presentation, path, kwargs.get("max_slides"))
        return self._write(Presentation, path, kwargs.get("slides"))

    def _read(self, Presentation, path: str, max_slides) -> str:
        if not os.path.exists(path):
            return f"Error: file not found: {path}"
        try:
            prs = Presentation(path)
        except Exception as e:
            return f"Error: failed to open pptx: {e}"

        limit = max_slides if isinstance(max_slides, int) and max_slides > 0 else 50
        total_slides = len(prs.slides)
        # python-pptx 的 Slides 不支持切片，手动遍历到 limit
        parts: List[str] = []
        for i in range(min(limit, total_slides)):
            slide = prs.slides[i]
            title = ""
            texts: List[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                shape_text = "\n".join(p.text for p in tf.paragraphs)
                if shape == slide.shapes.title:
                    title = shape_text
                else:
                    if shape_text:
                        texts.append(shape_text)
            block = f"[Slide {i + 1}] Title: {title}"
            if texts:
                block += "\n" + "\n".join(texts)
            parts.append(block)
        if total_slides > limit:
            parts.append(f"... [已截断，共 {total_slides} 张幻灯片，仅显示前 {limit} 张]")
        return "\n\n".join(parts)

    def _write(self, Presentation, path: str, slides) -> str:
        if slides is None:
            return "Error: slides is required for write action"
        try:
            parsed = _json.loads(slides) if isinstance(slides, str) else slides
        except Exception as e:
            return f"Error: failed to parse slides as JSON: {e}"

        if not isinstance(parsed, list):
            return "Error: slides must be a list of objects"

        parent_dir = os.path.dirname(path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)

        try:
            prs = Presentation()
            # 默认模板：
            #   slide_layouts[0] = Title Slide（仅标题）
            #   slide_layouts[1] = Title and Content（标题+正文）
            for item in parsed:
                if not isinstance(item, dict):
                    return "Error: each slide item must be an object"
                title = str(item.get("title", ""))
                content = str(item.get("content", ""))

                # 每项创建一张幻灯片，使用 layout[1]（Title and Content）
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                if title_shape is not None:
                    title_shape.text = title
                if content:
                    # placeholder idx 1 通常是正文 body
                    for ph in slide.placeholders:
                        try:
                            if ph.placeholder_format.idx == 1:
                                ph.text = content
                                break
                        except Exception:
                            continue
            prs.save(path)
        except Exception as e:
            return f"Error: failed to write pptx: {e}"
        return f"PowerPoint 文档已成功写入：{path}"
